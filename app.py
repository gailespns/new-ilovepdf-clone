"""
PDF Tools Backend — iLovePDF Clone (9 Features)
================================================
Original 5:
  1. Merge PDF         — Combine multiple PDFs into one
  2. Split PDF         — Split into pages or ranges
  3. Compress PDF      — Reduce size via Ghostscript / pikepdf fallback
  4. PDF → Word        — Convert to .docx via pdf2docx
  5. PDF → PowerPoint  — Rasterise pages into .pptx via pdf2image + python-pptx

New 4 (from screenshot):
  6. PDF → Excel       — Extract tables from PDF into .xlsx via pdfplumber + openpyxl
  7. Word → PDF        — Convert .docx/.doc to PDF via LibreOffice
  8. PowerPoint → PDF  — Convert .pptx/.ppt to PDF via LibreOffice
  9. Excel → PDF       — Convert .xlsx/.xls to PDF via LibreOffice
"""

import os, io, uuid, zipfile, shutil, subprocess, tempfile, logging
from pathlib import Path
from flask import Flask, request, jsonify, send_file, after_this_request
from pypdf import PdfReader, PdfWriter
import pikepdf
from pdf2docx import Converter as PdfToDocxConverter
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import pdfplumber
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
import cv2
import numpy as np

logging.basicConfig(level=logging.INFO, format="%(levelname)s  %(message)s")
log = logging.getLogger(__name__)

app = Flask(__name__, static_folder="static", static_url_path="")
UPLOAD_DIR = Path(tempfile.gettempdir()) / "pdf_tools_uploads"
OUTPUT_DIR = Path(tempfile.gettempdir()) / "pdf_tools_outputs"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)
app.config["MAX_CONTENT_LENGTH"] = 500 * 1024 * 1024  # 500 MB

PDF_EXT   = {".pdf"}
WORD_EXT  = {".doc", ".docx"}
PPTX_EXT  = {".ppt", ".pptx"}
EXCEL_EXT = {".xls", ".xlsx"}
LIBREOFFICE = shutil.which("libreoffice") or shutil.which("soffice")

# ── Helpers ───────────────────────────────────────────────────────────────────

def _uid(): return uuid.uuid4().hex

def _save_upload(fs, suffix=".pdf"):
    dest = UPLOAD_DIR / f"{_uid()}{suffix}"
    fs.save(str(dest))
    log.info("Saved upload → %s  (%.1f KB)", dest.name, dest.stat().st_size / 1024)
    return dest

def _out(suffix): return OUTPUT_DIR / f"{_uid()}{suffix}"

def _send(path, download_name, mimetype="application/octet-stream"):
    @after_this_request
    def _cleanup(r):
        try: path.unlink(missing_ok=True)
        except: pass
        return r
    return send_file(str(path), as_attachment=True, download_name=download_name, mimetype=mimetype)

def _pdf_page_count(path):
    return len(PdfReader(str(path)).pages)

def _parse_page_ranges(s, total):
    s = (s or "").strip().lower()

    if not s or s == "each":
        return [[i] for i in range(total)]
    if s == "all":
        return [list(range(total))]

    groups = []
    for part in s.split(","):
        part = part.strip()
        if not part:
            continue
        try:
            if "-" in part:
                lo, hi = part.split("-", 1)
                lo = int(lo.strip())
                hi = int(hi.strip())
                lo = max(1, lo)
                hi = min(total, hi)
                if lo <= hi:
                    groups.append(list(range(lo - 1, hi)))
            else:
                idx = int(part.strip())
                if 1 <= idx <= total:
                    groups.append([idx - 1])
        except Exception:
            log.warning("Invalid range ignored: %s", part)

    return groups

def _libreoffice_convert(input_path, fmt):
    """Run LibreOffice headless conversion. Returns output Path or None."""
    if not LIBREOFFICE: return None
    work = OUTPUT_DIR / _uid()
    work.mkdir()
    src = work / input_path.name
    shutil.copy(input_path, src)
    cmd = [LIBREOFFICE, "--headless", "--norestore",
           "--convert-to", fmt, "--outdir", str(work), str(src)]
    log.info("LibreOffice: %s", " ".join(cmd))
    res = subprocess.run(cmd, capture_output=True, timeout=120)
    if res.returncode != 0:
        log.error("LO stderr: %s", res.stderr.decode())
        shutil.rmtree(work, ignore_errors=True)
        return None
    outputs = [p for p in work.iterdir() if p.suffix.lower() == f".{fmt}"]
    if not outputs:
        shutil.rmtree(work, ignore_errors=True)
        return None
    out = _out(f".{fmt}")
    shutil.move(str(outputs[0]), out)
    shutil.rmtree(work, ignore_errors=True)
    log.info("  LO output %.1f KB", out.stat().st_size / 1024)
    return out

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 1 — MERGE PDF
# ═══════════════════════════════════════════════════════════════════════════

@app.route("/api/merge", methods=["POST"])
def merge_pdf():
    """
    POST /api/merge · field: files[] (2+ PDFs)
    Returns: merged.pdf

    Algorithm: Open each uploaded PDF with PdfReader; append every page
    sequentially to a shared PdfWriter. Write to temp file, stream back.
    All temp inputs deleted in finally block.
    """
    files = request.files.getlist("files[]")
    if len(files) < 2:
        return jsonify(error="Please upload at least 2 PDF files."), 400
    saved = []
    try:
        writer = PdfWriter()
        for f in files:
            if Path(f.filename).suffix.lower() not in PDF_EXT:
                return jsonify(error=f"'{f.filename}' is not a PDF."), 400
            p = _save_upload(f)
            saved.append(p)
            reader = PdfReader(str(p))
            log.info("Merging '%s' → %d pages", f.filename, len(reader.pages))
            for page in reader.pages:
                writer.add_page(page)
        out = _out(".pdf")
        with open(out, "wb") as fh: writer.write(fh)
        log.info("Merge done: %d files, %.1f KB", len(files), out.stat().st_size/1024)
        return _send(out, "merged.pdf", "application/pdf")
    finally:
        for p in saved: p.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 2 — SPLIT PDF
# ═══════════════════════════════════════════════════════════════════════════

@app.route("/api/split", methods=["POST"])
def split_pdf():
    f = request.files.get("file")
    if not f:
        return jsonify(error="No file provided."), 400
    if Path(f.filename).suffix.lower() not in PDF_EXT:
        return jsonify(error="File must be a PDF."), 400

    ranges_str = request.form.get("ranges", "each")
    input_path = _save_upload(f)
    output_paths = []

    try:
        reader = PdfReader(str(input_path))
        total = len(reader.pages)
        log.info("Split: '%s' has %d pages, ranges='%s'", f.filename, total, ranges_str)

        groups = _parse_page_ranges(ranges_str, total)
        if not groups:
            return jsonify(error="No valid page ranges found."), 400

        for idx, page_indices in enumerate(groups):
            writer = PdfWriter()
            for pi in page_indices:
                writer.add_page(reader.pages[pi])

            out = _out(".pdf")
            with open(out, "wb") as fh:
                writer.write(fh)
            output_paths.append(out)

            log.info(
                "Part %d/%d: pages %s",
                idx + 1,
                len(groups),
                [p + 1 for p in page_indices]
            )

        if len(output_paths) == 1:
            stem = Path(f.filename).stem
            return _send(output_paths[0], f"{stem}_split.pdf", "application/pdf")

        zip_path = _out(".zip")
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, p in enumerate(output_paths, 1):
                zf.write(p, f"page_{i:03d}.pdf")

        @after_this_request
        def _cleanup(response):
            for p in output_paths:
                try:
                    p.unlink(missing_ok=True)
                except Exception:
                    pass
            try:
                zip_path.unlink(missing_ok=True)
            except Exception:
                pass
            return response

        return send_file(
            str(zip_path),
            as_attachment=True,
            download_name="split_pages.zip",
            mimetype="application/zip",
        )

    except Exception as e:
        log.exception("Split error")
        return jsonify(error=str(e)), 500

    finally:
        input_path.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 3 — COMPRESS PDF
# ═══════════════════════════════════════════════════════════════════════════

GS_SETTINGS = {"screen":"/screen","ebook":"/ebook","printer":"/printer","prepress":"/prepress"}

def _compress_gs(src, dst, quality):
    gs = shutil.which("gs") or shutil.which("gswin64c")
    if not gs: return False
    cmd = [gs,"-sDEVICE=pdfwrite","-dCompatibilityLevel=1.4",
           f"-dPDFSETTINGS={GS_SETTINGS.get(quality,'/ebook')}",
           "-dNOPAUSE","-dQUIET","-dBATCH",f"-sOutputFile={dst}",str(src)]
    log.info("GS: %s", " ".join(cmd))
    res = subprocess.run(cmd, capture_output=True, timeout=120)
    return res.returncode == 0

def _compress_pikepdf(src, dst):
    with pikepdf.open(str(src)) as pdf:
        pdf.save(str(dst), compress_streams=True, preserve_pdfa=False,
                 object_stream_mode=pikepdf.ObjectStreamMode.generate)

@app.route("/api/compress", methods=["POST"])
def compress_pdf():
    """
    POST /api/compress · fields: file, quality (screen|ebook|printer|prepress)
    Returns: compressed_<n>.pdf

    Algorithm: Primary → Ghostscript -dPDFSETTINGS (lossy image downsampling
    + stream re-encoding). Fallback → pikepdf lossless (dedup + repack).
    Guard: if output >= input, return original unchanged.

    Quality presets:
      screen   /screen   ~72 dpi  — smallest, email/web
      ebook    /ebook    ~150 dpi — balanced (default)
      printer  /printer  ~300 dpi — high quality print
      prepress /prepress ~300 dpi + colour preservation
    """
    f = request.files.get("file")
    if not f: return jsonify(error="No file provided."), 400
    if Path(f.filename).suffix.lower() not in PDF_EXT:
        return jsonify(error="File must be a PDF."), 400
    quality = request.form.get("quality", "ebook")
    if quality not in GS_SETTINGS:
        return jsonify(error=f"Invalid quality. Choose: {list(GS_SETTINGS)}"), 400
    input_path = _save_upload(f)
    out_path = _out(".pdf")
    try:
        orig = input_path.stat().st_size
        log.info("Compress '%s' %.1f KB quality=%s", f.filename, orig/1024, quality)
        ok = _compress_gs(input_path, out_path, quality)
        method = "ghostscript"
        if not ok or not out_path.exists():
            log.warning("GS failed — pikepdf fallback")
            _compress_pikepdf(input_path, out_path)
            method = "pikepdf"
        comp = out_path.stat().st_size
        log.info("  %s %.1f KB  %.1f%% reduction", method, comp/1024, (1-comp/orig)*100)
        if comp >= orig:
            out_path.unlink(missing_ok=True)
            shutil.copy(input_path, out_path)
        stem = Path(f.filename).stem
        return _send(out_path, f"compressed_{stem}.pdf", "application/pdf")
    finally:
        input_path.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 4 — PDF TO WORD
# ═══════════════════════════════════════════════════════════════════════════

@app.route("/api/pdf-to-word", methods=["POST"])
def pdf_to_word():
    """
    POST /api/pdf-to-word · fields: file, start_page, end_page
    Returns: <n>.docx

    Algorithm: pdf2docx.Converter renders pages via PyMuPDF at high DPI,
    then a rule-based layout analyser detects text blocks / fonts / images /
    tables / columns and writes python-docx XML mirroring the layout.
    Page range maps to 0-based start/end args in Converter.convert().
    """
    f = request.files.get("file")
    if not f: return jsonify(error="No file provided."), 400
    if Path(f.filename).suffix.lower() not in PDF_EXT:
        return jsonify(error="File must be a PDF."), 400
    input_path = _save_upload(f)
    out_path = _out(".docx")
    try:
        total = _pdf_page_count(input_path)
        start = max(0, int(request.form.get("start_page", 1)) - 1)
        end   = min(total-1, int(request.form.get("end_page", total)) - 1)
        log.info("PDF→Word '%s' pages %d-%d of %d", f.filename, start+1, end+1, total)
        cv = PdfToDocxConverter(str(input_path))
        cv.convert(str(out_path), start=start, end=end)
        cv.close()
        log.info("  %.1f KB", out_path.stat().st_size/1024)
        stem = Path(f.filename).stem
        return _send(out_path, f"{stem}.docx",
                     "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    finally:
        input_path.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 5 — PDF TO POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════

SLIDE_W, SLIDE_H, DEFAULT_DPI = 10.0, 7.5, 150

@app.route("/api/pdf-to-pptx", methods=["POST"])
def pdf_to_pptx():
    """
    POST /api/pdf-to-pptx · fields: file, dpi (72-300, default 150)
    Returns: <n>.pptx (one slide per PDF page)

    Algorithm:
      1. pdf2image.convert_from_path() → poppler pdftoppm → PIL Images at DPI
      2. Each Image saved as temp PNG
      3. python-pptx Presentation (10×7.5 in, blank layout) built; each PNG
         inserted as a full-slide picture (left=0,top=0,w=slide_w,h=slide_h)
      4. Temp PNGs deleted after PPTX written
    Higher DPI → crisper text, larger file, slower conversion.
    """
    f = request.files.get("file")
    if not f: return jsonify(error="No file provided."), 400
    if Path(f.filename).suffix.lower() not in PDF_EXT:
        return jsonify(error="File must be a PDF."), 400
    try: dpi = max(72, min(300, int(request.form.get("dpi", DEFAULT_DPI))))
    except ValueError: return jsonify(error="Invalid DPI."), 400
    input_path = _save_upload(f)
    temp_imgs = []
    out_path = _out(".pptx")
    try:
        total = _pdf_page_count(input_path)
        log.info("PDF→PPTX '%s' %d pages dpi=%d", f.filename, total, dpi)
        pil_imgs = convert_from_path(str(input_path), dpi=dpi, fmt="png", thread_count=2)
        for i, img in enumerate(pil_imgs):
            tmp = OUTPUT_DIR / f"{_uid()}_p{i}.png"
            img.save(str(tmp), "PNG")
            temp_imgs.append(tmp)
        prs = Presentation()
        prs.slide_width  = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        blank = prs.slide_layouts[6]
        for ip in temp_imgs:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(ip), Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        prs.save(str(out_path))
        log.info("  %.1f KB", out_path.stat().st_size/1024)
        stem = Path(f.filename).stem
        return _send(out_path, f"{stem}.pptx",
                     "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    finally:
        input_path.unlink(missing_ok=True)
        for p in temp_imgs: p.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 6 — PDF TO EXCEL
# ═══════════════════════════════════════════════════════════════════════════

def _style_header(ws, row, ncols):
    fill = PatternFill("solid", fgColor="2E4057")
    font = Font(bold=True, color="FFFFFF", size=11)
    thin = Side(style="thin", color="CCCCCC")
    bdr  = Border(left=thin, right=thin, top=thin, bottom=thin)
    for c in range(1, ncols+1):
        cell = ws.cell(row=row, column=c)
        cell.fill, cell.font, cell.border = fill, font, bdr
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

def _style_data(ws, row, ncols, alt):
    fill = PatternFill("solid", fgColor="F2F4F7" if alt else "FFFFFF")
    thin = Side(style="thin", color="DDDDDD")
    bdr  = Border(left=thin, right=thin, top=thin, bottom=thin)
    for c in range(1, ncols+1):
        cell = ws.cell(row=row, column=c)
        cell.fill, cell.border = fill, bdr
        cell.alignment = Alignment(wrap_text=True, vertical="top")

def _auto_width(ws):
    for col in ws.columns:
        ltr = get_column_letter(col[0].column)
        w = max((len(str(c.value or "")) for c in col), default=0)
        ws.column_dimensions[ltr].width = min(w + 4, 62)

@app.route("/api/pdf-to-excel", methods=["POST"])
@app.route("/api/pdf-to-excel", methods=["POST"])
def pdf_to_excel():
    f = request.files.get("file")
    if not f:
        return jsonify(error="No file provided."), 400
    if Path(f.filename).suffix.lower() not in PDF_EXT:
        return jsonify(error="File must be a PDF."), 400

    input_path = _save_upload(f)
    out_path = _out(".xlsx")

    def looks_garbled(text: str) -> bool:
        if not text:
            return True
        bad = sum(1 for ch in text if not (ch.isalnum() or ch.isspace() or ch in ".,:/()-_%'\"&"))
        return (bad / max(len(text), 1)) > 0.12

    def clean_cell(val):
        if val is None:
            return ""
        return " ".join(str(val).replace("\n", " ").split())

    try:
        log.info("PDF→Excel '%s'", f.filename)
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        total_tables = 0

        with pdfplumber.open(str(input_path)) as pdf:
            total_pages = len(pdf.pages)

            for pg_idx, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables({
                    "vertical_strategy": "lines",
                    "horizontal_strategy": "lines",
                    "snap_tolerance": 3,
                    "join_tolerance": 3,
                    "edge_min_length": 3,
                    "min_words_vertical": 1,
                    "min_words_horizontal": 1,
                })

                if not tables:
                    tables = page.extract_tables({
                        "vertical_strategy": "text",
                        "horizontal_strategy": "text",
                        "snap_tolerance": 3,
                        "join_tolerance": 3,
                        "min_words_vertical": 2,
                        "min_words_horizontal": 1,
                    })

                usable_tables = []
                for table in tables or []:
                    cleaned = []
                    sample_text = []
                    for row in table:
                        if not row:
                            continue
                        cleaned_row = [clean_cell(cell) for cell in row]
                        cleaned.append(cleaned_row)
                        sample_text.extend(cleaned_row)

                    joined = " ".join(sample_text[:20])
                    if cleaned and not looks_garbled(joined):
                        usable_tables.append(cleaned)

                if usable_tables:
                    log.info("  Page %d/%d: %d usable table(s)", pg_idx, total_pages, len(usable_tables))
                    for tbl_idx, table in enumerate(usable_tables, 1):
                        total_tables += 1
                        ws = wb.create_sheet(title=f"P{pg_idx}-T{tbl_idx}")
                        ncols = max(len(r) for r in table)

                        for rn, row in enumerate(table, 1):
                            padded = list(row) + [""] * (ncols - len(row))
                            for cn, val in enumerate(padded, 1):
                                ws.cell(row=rn, column=cn, value=val)

                            if rn == 1:
                                _style_header(ws, rn, ncols)
                            else:
                                _style_data(ws, rn, ncols, rn % 2 == 0)

                        ws.freeze_panes = "A2"
                        _auto_width(ws)
                    continue

                # OCR fallback
                log.info("  Page %d/%d: no usable tables, using OCR fallback", pg_idx, total_pages)
                images = convert_from_path(
                    str(input_path),
                    dpi=250,
                    first_page=pg_idx,
                    last_page=pg_idx,
                    fmt="png"
                )
                if not images:
                    continue

                img = images[0]
                img_np = np.array(img)
                gray = cv2.cvtColor(img_np, cv2.COLOR_RGB2GRAY)
                gray = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)[1]

                ocr_text = pytesseract.image_to_string(
                    gray,
                    config="--oem 3 --psm 6"
                )
                lines = [line.strip() for line in ocr_text.splitlines() if line.strip()]

                if not lines:
                    continue

                total_tables += 1
                ws = wb.create_sheet(title=f"P{pg_idx}-OCR")
                ws.cell(row=1, column=1, value="OCR Extracted Text")
                _style_header(ws, 1, 1)

                rn = 2
                for line in lines:
                    ws.cell(row=rn, column=1, value=" ".join(line.split()))
                    _style_data(ws, rn, 1, rn % 2 == 0)
                    rn += 1

                ws.freeze_panes = "A2"
                ws.column_dimensions["A"].width = 120

            if total_tables == 0:
                return jsonify(error="No readable tables or OCR text detected in this PDF."), 400

        wb.save(str(out_path))
        stem = Path(f.filename).stem
        return _send(
            out_path,
            f"{stem}.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    finally:
        input_path.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 7 — WORD TO PDF
# ═══════════════════════════════════════════════════════════════════════════

@app.route("/api/word-to-pdf", methods=["POST"])
def word_to_pdf():
    """
    POST /api/word-to-pdf · field: file (.doc / .docx)
    Returns: <n>.pdf

    Algorithm: LibreOffice headless mode:
      libreoffice --headless --convert-to pdf --outdir <dir> <file>
    LibreOffice's ODF/OOXML rendering engine reproduces fonts, images,
    tables, headers/footers, and styles. Output is fully searchable PDF
    (not a raster). Source file copied to a dedicated temp dir so LibreOffice
    --outdir places output in a known location.
    """
    f = request.files.get("file")
    if not f: return jsonify(error="No file provided."), 400
    ext = Path(f.filename).suffix.lower()
    if ext not in WORD_EXT: return jsonify(error="File must be .doc or .docx"), 400
    if not LIBREOFFICE: return jsonify(error="LibreOffice not available."), 500
    input_path = _save_upload(f, suffix=ext)
    try:
        log.info("Word→PDF '%s'", f.filename)
        out = _libreoffice_convert(input_path, "pdf")
        if not out: return jsonify(error="Conversion failed."), 500
        return _send(out, f"{Path(f.filename).stem}.pdf", "application/pdf")
    finally:
        input_path.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 8 — POWERPOINT TO PDF
# ═══════════════════════════════════════════════════════════════════════════

@app.route("/api/pptx-to-pdf", methods=["POST"])
def pptx_to_pdf():
    """
    POST /api/pptx-to-pdf · field: file (.ppt / .pptx)
    Returns: <n>.pdf

    Algorithm: Same LibreOffice headless pipeline as Word→PDF.
    LibreOffice Impress renders each slide to a PDF page. Slide transitions
    and animations are not preserved (PDF is static), but all text, images,
    shapes, and charts are rendered. One PDF page per slide.
    """
    f = request.files.get("file")
    if not f: return jsonify(error="No file provided."), 400
    ext = Path(f.filename).suffix.lower()
    if ext not in PPTX_EXT: return jsonify(error="File must be .ppt or .pptx"), 400
    if not LIBREOFFICE: return jsonify(error="LibreOffice not available."), 500
    input_path = _save_upload(f, suffix=ext)
    try:
        log.info("PPTX→PDF '%s'", f.filename)
        out = _libreoffice_convert(input_path, "pdf")
        if not out: return jsonify(error="Conversion failed."), 500
        return _send(out, f"{Path(f.filename).stem}.pdf", "application/pdf")
    finally:
        input_path.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# FEATURE 9 — EXCEL TO PDF
# ═══════════════════════════════════════════════════════════════════════════

@app.route("/api/excel-to-pdf", methods=["POST"])
def excel_to_pdf():
    """
    POST /api/excel-to-pdf · field: file (.xls / .xlsx)
    Returns: <n>.pdf

    Algorithm: Same LibreOffice headless pipeline. LibreOffice Calc renders
    each non-empty worksheet to one or more PDF pages using the workbook's
    embedded print-area settings (or the full used range if none are set).
    All sheets in the workbook are exported into one PDF.
    """
    f = request.files.get("file")
    if not f: return jsonify(error="No file provided."), 400
    ext = Path(f.filename).suffix.lower()
    if ext not in EXCEL_EXT: return jsonify(error="File must be .xls or .xlsx"), 400
    if not LIBREOFFICE: return jsonify(error="LibreOffice not available."), 500
    input_path = _save_upload(f, suffix=ext)
    try:
        log.info("Excel→PDF '%s'", f.filename)
        out = _libreoffice_convert(input_path, "pdf")
        if not out: return jsonify(error="Conversion failed."), 500
        return _send(out, f"{Path(f.filename).stem}.pdf", "application/pdf")
    finally:
        input_path.unlink(missing_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# UTILITY ENDPOINTS
# ═══════════════════════════════════════════════════════════════════════════

@app.route("/api/info", methods=["POST"])
def pdf_info():
    f = request.files.get("file")
    if not f: return jsonify(error="No file provided."), 400
    p = _save_upload(f)
    try:
        reader = PdfReader(str(p))
        meta = reader.metadata or {}
        return jsonify(filename=f.filename, pages=len(reader.pages),
                       file_size_kb=round(p.stat().st_size/1024, 1),
                       title=meta.get("/Title",""), author=meta.get("/Author",""),
                       creator=meta.get("/Creator",""), encrypted=reader.is_encrypted)
    finally:
        p.unlink(missing_ok=True)

@app.route("/api/health")
def health():
    return jsonify(status="ok",
                   ghostscript=bool(shutil.which("gs") or shutil.which("gswin64c")),
                   pdftoppm=bool(shutil.which("pdftoppm")),
                   libreoffice=bool(LIBREOFFICE))

@app.route("/")
def index():
    return app.send_static_file("index.html")

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(debug=False, host="0.0.0.0", port=port)